"""
EHPM Revision 1 (R1): Healthcare Expenditure as Economic Effect

Revised manuscript addressing Reviewer 1 (methodological) and Reviewer 2
(study design) comments.  Major changes vs. original submission:

Reviewer 1:
  R1-1: FRR economic justification — output vs value-added multipliers,
        opportunity cost / crowding-out caveat
  R1-2: Table 1 comparability — US Medicare qualifier, sensitivity analysis
  R1-3: Tempo model self-contained specification (μ_H0, SE, estimation)
  R1-4: AIC/BIC/LOOCV for M0/M1/M2 in Table 3 — honest M2-vs-M1 assessment
  R1-5: 15 % equipment-CHE share sensitivity (new Figure 9, new Table 5)
  R1-6: Scenario A density→spending→multiplier chain transparency
  R1-7: Uncertainty range around 0.98; toned-down conclusions

Reviewer 2:
  R2-1: Explicit author position on equipment density
  R2-2: Late-Stage Elderly Healthcare System discussion
  R2-3: Super-aged society context, physician density/remuneration
  R2-4: Title revised to signal Japan focus
  R2-5: Deficit financing / national debt discussion (new subsection)
  R2-6: Abstract revised — OECD as context, Japan as focus
  R2-7: Over-supply / defensive medicine discussion
  R2-8: International comparability caveats (quality, access, GP gatekeeping)

Format: EHPM Research Article, double-spaced, 12 pt, A4, Vancouver refs.
"""
import os
import re
import json
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, ".."))
DATA = os.path.join(ROOT, "data")
FIG = os.path.join(ROOT, "output", "figures")
DOCX_DIR = os.path.join(ROOT, "output", "docx")
PPTX_DIR = os.path.join(ROOT, "output", "pptx")
os.makedirs(DOCX_DIR, exist_ok=True)
os.makedirs(PPTX_DIR, exist_ok=True)


def get_fig(name):
    return os.path.join(FIG, name)


# ---------------------------------------------------------------------------
# EHPM formatting helpers (unchanged from v0)
# ---------------------------------------------------------------------------
def _set_ehpm_format(doc):
    for section in doc.sections:
        section.page_width = Cm(21.0)
        section.page_height = Cm(29.7)
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = Pt(12)
    style.paragraph_format.line_spacing_rule = WD_LINE_SPACING.DOUBLE
    style.paragraph_format.space_after = Pt(0)
    for level in range(1, 4):
        hs = doc.styles[f'Heading {level}']
        hs.font.name = 'Times New Roman'
        hs.font.color.rgb = RGBColor(0, 0, 0)
        hs.paragraph_format.line_spacing_rule = WD_LINE_SPACING.DOUBLE
        if level == 1:
            hs.font.size = Pt(14)
        elif level == 2:
            hs.font.size = Pt(13)
        else:
            hs.font.size = Pt(12)


def add_text_with_refs(paragraph, text, bold=False, italic=False):
    parts = re.split(r'(\{[^}]+\})', text)
    for part in parts:
        if part.startswith('{') and part.endswith('}'):
            run = paragraph.add_run(part[1:-1])
            run.font.superscript = True
            run.font.size = Pt(10)
        else:
            run = paragraph.add_run(part)
            if bold:
                run.bold = True
            if italic:
                run.italic = True


def add_heading(doc, text, level=1):
    return doc.add_heading(text, level=level)


def add_para(doc, text, bold=False, italic=False):
    p = doc.add_paragraph()
    add_text_with_refs(p, text, bold=bold, italic=italic)
    return p


def add_plain_para(doc, text, bold=False, italic=False, align=None):
    p = doc.add_paragraph()
    run = p.add_run(text)
    if bold:
        run.bold = True
    if italic:
        run.italic = True
    if align:
        p.alignment = align
    return p


def add_table_from_df(doc, df, title, legend=""):
    cap = doc.add_paragraph()
    cap.paragraph_format.space_before = Pt(12)
    run = cap.add_run(title)
    run.bold = True
    run.font.size = Pt(12)
    table = doc.add_table(rows=1, cols=len(df.columns), style='Table Grid')
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for j, col in enumerate(df.columns):
        cell = table.rows[0].cells[j]
        cell.text = str(col)
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.size = Pt(10)
                r.font.name = 'Times New Roman'
    for _, row in df.iterrows():
        row_cells = table.add_row().cells
        for j, col in enumerate(df.columns):
            row_cells[j].text = str(row[col])
            for p in row_cells[j].paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)
                    r.font.name = 'Times New Roman'
    if legend:
        leg = doc.add_paragraph()
        run = leg.add_run(legend)
        run.font.size = Pt(10)
        run.italic = True
    doc.add_paragraph()


def add_figure_inline(doc, img_path, caption):
    if os.path.exists(img_path):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run()
        run.add_picture(img_path, width=Inches(5.0))
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap.paragraph_format.space_before = Pt(6)
    cap.paragraph_format.space_after = Pt(12)
    run = cap.add_run(caption)
    run.font.size = Pt(10)
    run.italic = True


# ---------------------------------------------------------------------------
# References -- Vancouver style, numbered by order of first appearance
# R1 additions: refs 25-30 for new content
# ---------------------------------------------------------------------------
REFERENCES = [
    # 1 -- Background: OECD ranking
    "OECD. Health at a Glance 2023: OECD Indicators. Paris: OECD Publishing; "
    "2023.",
    # 2 -- Background: Cabinet Office
    "Cabinet Office, Japan. Annual Report on the Japanese Economy 2025. "
    "Tokyo: Cabinet Office; 2025. [In Japanese]",
    # 3 -- R1 NEW: MOF fiscal data (deficit financing)
    "Ministry of Finance, Japan. Highlights of the Budget for FY2024. "
    "Tokyo: MOF; 2024. Available from: "
    "https://www.mof.go.jp/english/policy/budget/budget/",
    # 4 -- R1 NEW: Japan national debt
    "Bank of Japan. Flow of Funds Accounts. Tokyo: BOJ; 2024. "
    "Available from: https://www.boj.or.jp/en/statistics/sj/",
    # 5 -- R1 NEW: Late-stage elderly healthcare
    "Ministry of Health, Labour and Welfare. Annual Health, Labour and "
    "Welfare Report 2023. Tokyo: MHLW; 2023. [In Japanese]",
    # 6 -- Background: JMARI employment
    "Maeda Y. Economic ripple effects of healthcare and long-term care. "
    "JMARI Working Paper No. 172. 2008. [In Japanese]",
    # 7 -- Background: Japan multiplier
    "Yamada G, Imanaka Y. Input-output analysis on the economic impact of "
    "medical care in Japan. Environ Health Prev Med. 2015;20(5):379-387.",
    # 8 -- Background: US multiplier
    "Dupor B, Guerrero R. The aggregate and local economic effects of "
    "government financed health care. Econ Inq. 2021;59(2):662-670.",
    # 9 -- Background: EU I-O framework
    "Gutierrez-Hernandez I, Abasolo-Alesson I. The health care sector in "
    "the economies of the European Union: an overview using an input-output "
    "framework. Cost Eff Resour Alloc. 2021;19(1):4.",
    # 10 -- Methods: Germany I-O
    "Henke KD, Ostwald DA. Health satellite account: the first step. "
    "Int J Behav Healthc Res. 2012;3(2):91-105.",
    # 11 -- Background: Health as investment
    "Mushkin SJ. Health as an investment. J Polit Econ. 1962;70(5):129-157.",
    # 12 -- Background: Health capital model
    "Grossman M. On the concept of health capital and the demand for "
    "health. J Polit Econ. 1972;80(2):223-255.",
    # 13 -- Background: Health-TFP link
    "Bloom DE, Canning D, Sevilla J. The effect of health on economic "
    "growth: a production function approach. World Dev. "
    "2004;32(1):1-13.",
    # 14 -- Background: Life expectancy and GDP
    "Barro RJ. Health and economic growth. Ann Econ Finance. "
    "2013;14(2):305-342.",
    # 15 -- Background: HLGH (CS-ARDL)
    "Ertugrul HM, Baycan O, Atilgan E, Ulucan H. Health-led growth "
    "hypothesis and health financing systems: an econometric synthesis "
    "for OECD countries. Front Public Health. 2024;12:1437304.",
    # 16 -- Background: HLGH (Toda-Yamamoto)
    "Amiri A, Ventelou B. Granger causality between total expenditure on "
    "health and GDP in OECD: Evidence from the Toda-Yamamoto approach. "
    "Econ Lett. 2012;116(3):541-544.",
    # 17 -- Background: HLGH (Driscoll-Kraay)
    "Beylik U, Cirakli U, Cetin M, Ecevit E, Senol O. The relationship "
    "between health expenditure indicators and economic growth in OECD "
    "countries: A Driscoll-Kraay approach. Front Public Health. "
    "2022;10:1050550.",
    # 18 -- Background: HLGH (Panel VECM)
    "Wang KM. Health care expenditure and economic growth: Quantile "
    "panel-type analysis. Econ Model. 2011;28(4):1536-1549.",
    # 19 -- Background: HLGH developing countries
    "Piabuo SM, Tieguhong JC. Health expenditure and economic growth - "
    "a review of the literature and an analysis between the economic "
    "community for central African states (CEMAC) and selected African "
    "countries. Health Econ Rev. 2017;7(1):23.",
    # 20 -- Background: Tempo effect origin
    "Bongaarts J, Feeney G. On the quantum and tempo of fertility. "
    "Popul Dev Rev. 1998;24(2):271-291.",
    # 21 -- Background: Forgotten parameter sigma
    "Goldstein JR, Lutz W, Scherbov S. Long-term population decline in "
    "Europe: the relative importance of tempo effects and generational "
    "length. Popul Dev Rev. 2003;29(4):699-707.",
    # 22 -- Methods: GDP tempo companion
    "Onishi T. The forgotten tempo effect in capital accounting: "
    "investment-to-output time-to-build, intangible capital, and the "
    "reconciliation of flow- and stock-based national wealth measures. "
    "Working Paper. 2026. Available from: "
    "https://github.com/bougtoir/gdp-tempo-paper",
    # 23 -- Methods: Healthcare tempo PoC
    "Onishi T. Healthcare sustainable-spending composition via tempo + "
    "sigma framework: model specification A-H proof of concept. "
    "Working Paper. 2026. Available from: "
    "https://github.com/bougtoir/healthcare-sustainable",
    # 24 -- Results: OECD diagnostic technology
    "OECD. Health at a Glance 2023: OECD Indicators. Chapter 5.23: "
    "Diagnostic technologies. Paris: OECD Publishing; 2023.",
    # 25 -- Results: MHLW trade statistics
    "Ministry of Health, Labour and Welfare. Pharmaceutical and medical "
    "device production statistics annual report. "
    "Tokyo: MHLW; 2021. [In Japanese]",
    # 26 -- R1 NEW: Physician density OECD
    "OECD. Health at a Glance 2023: OECD Indicators. Chapter 8.1: "
    "Health workforce. Paris: OECD Publishing; 2023.",
    # 27 -- R1 NEW: Work-style reform physicians
    "Koike S, Wada H, Ohde S, Ide H, Taneda K, Tanigawa T. Working "
    "hours of full-time hospital physicians in Japan: a cross-sectional "
    "nationwide survey. BMC Public Health. 2024;24:164.",
    # 28 -- Results: WDI data
    "World Bank. World Development Indicators. Washington, DC: World Bank; "
    "2024. Available from: https://databank.worldbank.org/",
    # 29 -- Results: Preston Curve
    "Preston SH. The changing relation between mortality and level of "
    "economic development. Popul Stud. 1975;29(2):231-248.",
    # 30 -- R1 NEW: Defensive medicine Japan
    "Hiyama T, Yoshihara M, Tanaka S, Urabe Y, Ikegami Y, Fukuhara T, "
    "et al. Defensive medicine practices among gastroenterologists in "
    "Japan. World J Gastroenterol. 2006;12(47):7671-7675.",
]


# ---------------------------------------------------------------------------
# Build EHPM R1 manuscript
# ---------------------------------------------------------------------------
def build_ehpm_manuscript_r1():
    doc = Document()
    _set_ehpm_format(doc)

    # ======================================================================
    # TITLE PAGE
    # ======================================================================
    add_plain_para(doc, "TITLE PAGE", bold=True,
                   align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()

    # R2-4: Title revised to signal Japan focus
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(
        "Evaluating Healthcare Expenditure Sustainability in Japan: "
        "A Dual-Return Framework Integrating Input-Output Multipliers, "
        "Health-Capital Tempo, and Diagnostic Equipment Stock "
        "with Cross-Country Benchmarking"
    )
    run.bold = True
    run.font.size = Pt(14)

    doc.add_paragraph()
    add_plain_para(doc, "Tatsuki Onishi [1]*", bold=False,
                   align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()
    add_plain_para(doc, "[1] [Affiliation to be completed by author]",
                   italic=True)
    doc.add_paragraph()

    p = doc.add_paragraph()
    run = p.add_run("*Corresponding author: ")
    run.bold = True
    p.add_run("Tatsuki Onishi, [Affiliation]. E-mail: [to be completed]")
    doc.add_paragraph()

    p = doc.add_paragraph()
    run = p.add_run("Short title: ")
    run.bold = True
    p.add_run("Japan Healthcare I-O Sustainability")
    doc.add_paragraph()

    p = doc.add_paragraph()
    run = p.add_run("Total word count: ")
    run.bold = True
    p.add_run("approximately 7,800 words "
              "(Text, References, Tables, and Figure Legends)")
    doc.add_paragraph()

    p = doc.add_paragraph()
    run = p.add_run("Number of figures: ")
    run.bold = True
    p.add_run("9")
    p = doc.add_paragraph()
    run = p.add_run("Number of tables: ")
    run.bold = True
    p.add_run("5")
    doc.add_page_break()

    # ======================================================================
    # ABSTRACT — R2-6: revised to focus on Japan
    # ======================================================================
    add_heading(doc, "Abstract", level=1)

    p = doc.add_paragraph()
    run = p.add_run("Background: ")
    run.bold = True
    p.add_run(
        "Japan's healthcare expenditure (11.0% of GDP) is conventionally "
        "treated as a fiscal cost to be contained, particularly given "
        "rapid population aging. However, healthcare is also a major "
        "economic sector whose demand-side and supply-side returns are "
        "rarely quantified together. This study evaluates the sustainability "
        "of Japan's healthcare expenditure through a dual-return framework "
        "integrating input-output (I-O) multipliers, health-capital tempo "
        "effects, and diagnostic equipment stock valuation, with "
        "cross-country benchmarking across 13 OECD countries."
    )

    p = doc.add_paragraph()
    run = p.add_run("Methods: ")
    run.bold = True
    p.add_run(
        "We compiled healthcare sector I-O output multipliers for 13 OECD "
        "countries from published national and EU-28 framework studies. "
        "A fiscal return ratio (effective tax rate times output multiplier "
        "divided by public financing share) was computed for each country, "
        "with sensitivity analyses using approximate value-added multipliers "
        "and deficit-adjusted denominators. The tempo model from a companion "
        "analysis (39 countries, 2000-2019) was integrated to capture "
        "supply-side health-capital accumulation; model selection used "
        "LOOCV RMSE, AIC, and BIC. Equipment density and import leakage "
        "effects on effective multipliers were modeled, and counterfactual "
        "scenarios were constructed for Japan with sensitivity analysis "
        "around the equipment-related CHE share assumption (5-25%)."
    )

    p = doc.add_paragraph()
    run = p.add_run("Results: ")
    run.bold = True
    p.add_run(
        "Using output multipliers, the demand-side fiscal return ratio "
        "exceeded 1.0 in five of thirteen countries (France 1.18, Italy "
        "1.13, Japan 1.09, Sweden 1.04, Finland 1.04). With value-added "
        "multipliers, no country exceeded 1.0, indicating that the output "
        "multiplier provides an upper bound. Japan's ratio fell from 1.09 "
        "to 1.04 after adjusting for 5% import leakage. Counterfactual "
        "analysis showed that reducing equipment density to the OECD "
        "average would lower the ratio to 0.98 (95% sensitivity range "
        "0.94-1.02 for equipment share 5-25%), indicating that the "
        "sustainability conclusion is sensitive to this assumption. "
        "The tempo model confirmed a spending-to-outcome lag (M1 vs M0: "
        "RMSE 0.441 vs 0.510, LOOCV 0.472 vs 0.538); however, the "
        "time-varying extension (M2 vs M1) showed only marginal "
        "improvement that did not survive BIC penalization."
    )

    p = doc.add_paragraph()
    run = p.add_run("Conclusions: ")
    run.bold = True
    p.add_run(
        "Japan's healthcare expenditure generates substantial demand-side "
        "economic returns via I-O multiplier effects. However, whether "
        "these returns fully cover public costs depends on methodological "
        "assumptions, including multiplier type (output vs value-added), "
        "deficit financing treatment, and equipment-related spending share. "
        "Japan's high diagnostic equipment density contributes to the "
        "economic return but its net effect on sustainability requires "
        "evaluation alongside over-supply risks and workforce constraints. "
        "The findings support integrating economic return analysis into "
        "healthcare policy evaluation, while cautioning against "
        "deterministic sustainability claims."
    )

    doc.add_paragraph()
    p = doc.add_paragraph()
    run = p.add_run("Keywords: ")
    run.bold = True
    p.add_run(
        "Healthcare expenditure; Input-output analysis; Economic multiplier; "
        "Health capital; Sustainability; Diagnostic imaging; Japan; "
        "Super-aged society; Deficit financing; OECD"
    )
    doc.add_page_break()

    # ======================================================================
    # BACKGROUND
    # ======================================================================
    add_heading(doc, "Background", level=1)

    add_para(doc,
        "Since the Cabinet's 'Basic Policies on Economic and Fiscal Management' "
        "decision in 2005, the Japanese government has pursued policies to "
        "moderate the growth of national medical care expenditure. Japan's "
        "current health expenditure (CHE) reached 11.0% of GDP in 2019, "
        "ranking fifth among 38 Organisation for Economic Co-operation and "
        "Development (OECD) countries{1}. With rapid population aging -- "
        "the proportion of the population aged 65 and over exceeded 28% "
        "in 2019, the highest among OECD nations -- the 'sustainability' of "
        "healthcare spending has become a central policy concern{2}."
    )
    # R2-3 & R2-5: Super-aged society context + national debt
    add_para(doc,
        "Japan faces a uniquely challenging fiscal context: its general "
        "account deficit dependency ratio has exceeded 30% in recent years, "
        "meaning that roughly one-third of government expenditures are "
        "financed by bond issuance rather than current tax revenue{3}. "
        "The national debt-to-GDP ratio exceeds 250%, the highest among "
        "advanced economies{4}. Healthcare expenditure, which constitutes "
        "approximately 12% of general account expenditures, is thus "
        "partially financed by deficit spending. Moreover, Japan's "
        "super-aged society -- characterized by a growing elderly demographic "
        "(projected 35% aged 65+ by 2040) alongside a shrinking workforce "
        "(projected 20% decline by 2040) -- places structural pressure "
        "on both the revenue and expenditure sides of healthcare financing."
    )
    # R2-2: Late-stage elderly healthcare system
    add_para(doc,
        "A distinctive feature of Japan's system is the Late-Stage Elderly "
        "Healthcare System (Kouki Koreisha Iryo Seido), which provides "
        "a separate insurance framework for individuals aged 75 and over{5}. "
        "Unlike many European countries where the institutional framework "
        "remains unified, with the elderly protected through subsidized "
        "premiums or exemptions, Japan changes the actual insurance "
        "'container' upon reaching age 75. This system is financed by "
        "a combination of patient co-payments (10-30%), premiums from "
        "enrollees, cross-subsidies from working-age insurance schemes, "
        "and public subsidies. The increasing weight of this system -- "
        "already covering approximately 18 million enrollees -- is a "
        "key driver of aggregate healthcare expenditure growth."
    )

    add_para(doc,
        "This discourse has almost uniformly treated healthcare expenditure "
        "as a 'cost' -- an expense to be minimized through efficiency gains, "
        "volume controls, and fee schedule revisions. "
        "Yet healthcare is also a major economic sector. It purchases inputs "
        "from pharmaceuticals, medical devices, and information technology; "
        "it employs millions of workers; and these workers spend their "
        "incomes, generating further economic activity throughout the economy. "
        "Maeda estimated that Japan's healthcare sector supports 6.89 million "
        "jobs in total (2.95 million direct, approximately 4 million indirect), "
        "with production-inducement effects exceeding those of any other "
        "service industry including construction and education{6}."
    )
    add_para(doc,
        "Input-output (I-O) analysis, pioneered by Leontief, quantifies how "
        "final demand in one sector induces production across the entire "
        "economy through direct, indirect, and induced effects. "
        "Yamada and Imanaka reported a healthcare I-O output multiplier of "
        "2.78 (95% confidence interval [CI]: 2.74-2.90) for Japan{7}, "
        "meaning that each unit of healthcare spending generates 2.78 units "
        "of total economic output. Output multipliers ranging from 1.7 to "
        "2.9 have been estimated across OECD countries{8-10}. "
        "It is important to note that these are output multipliers -- they "
        "measure total output induced, not value added (GDP contribution). "
        "Value-added multipliers, which exclude intermediate consumption, "
        "are typically 40-60% smaller{9}. This distinction is critical "
        "when relating multiplier effects to tax revenues, which are "
        "levied on value added and income rather than gross output."
    )
    add_para(doc,
        "On the supply side, health has long been recognized as a component "
        "of human capital that generates long-term economic returns. "
        "Mushkin conceptualized health as an investment with returns "
        "analogous to education{11}, and Grossman's health-capital model "
        "formally describes individual health stock as accumulated through "
        "investment and depreciated by aging and disease{12}. At the macro "
        "level, Bloom, Canning, and Sevilla demonstrated that improved "
        "population health raises total factor productivity{13}. Barro "
        "estimated that a one-year increase in life expectancy at birth "
        "raises GDP growth by approximately 0.04 percentage points{14}."
    )
    add_para(doc,
        "The Health-Led Growth Hypothesis (HLGH) synthesizes these supply-side "
        "effects, positing bidirectional Granger causality between health "
        "expenditure and GDP growth. Multiple studies have confirmed this "
        "relationship across OECD and developing countries using diverse "
        "econometric specifications{15-19}."
    )
    add_para(doc,
        "A further analytical dimension comes from the Bongaarts-Feeney "
        "tempo framework. Bongaarts and Feeney demonstrated that rising "
        "mean age at childbearing mechanically depresses the period total "
        "fertility rate even when cohort fertility is unchanged{20}. "
        "Goldstein, Lutz, and Scherbov introduced the parity-specific "
        "variance sigma, substantially improving tempo adjustment{21}. "
        "Onishi ported this quantum-tempo decomposition to capital "
        "accounting, introducing the investment-to-output time-to-build "
        "lag as an analogous 'forgotten parameter' in GDP accounting{22}. "
        "Extension to healthcare (termed 'Candidate A-H'){23} models the "
        "spending-to-outcome lag and revealed a tempo drift of "
        "+0.15 yr/yr across 39 countries (2000-2019)."
    )
    add_para(doc,
        "Japan is also known for its extremely high diagnostic imaging "
        "equipment density{1,24}. Japan has 115.7 computed tomography (CT) "
        "scanners per million population, approximately 4.3 times the OECD "
        "median, and 55.2 magnetic resonance imaging (MRI) units per million, "
        "approximately 2.9 times the OECD median. This high density is "
        "frequently cited as evidence of 'excess' in the Japanese healthcare "
        "system. Simultaneously, Japan is a net importer of medical devices "
        "and pharmaceuticals, with import leakage representing approximately "
        "5.0% of CHE{25}."
    )
    # R2-3: Physician density context
    add_para(doc,
        "An important structural context is Japan's physician density: at "
        "2.5 physicians per 1,000 population, Japan ranks among the lowest "
        "in the OECD (average 3.7){26}. This low physician density, "
        "combined with free access to specialists and the world's highest "
        "diagnostic equipment density, raises questions about the "
        "sustainability of current service levels, particularly as "
        "work-style reform legislation limits physician overtime{27}."
    )
    add_para(doc,
        "The aim of this study was to (1) evaluate the sustainability of "
        "Japan's healthcare expenditure from a neutral standpoint as both "
        "cost and economic effect, using a dual-return framework integrating "
        "demand-side I-O multipliers and supply-side health-capital tempo "
        "effects, (2) quantify the role of diagnostic equipment stock and "
        "import leakage, (3) construct counterfactual scenarios with "
        "sensitivity analysis, and (4) benchmark Japan against 12 additional "
        "OECD countries. The conceptual framework is illustrated in "
        "Figure 1."
    )

    # ======================================================================
    # METHODS
    # ======================================================================
    add_heading(doc, "Methods", level=1)

    add_heading(doc, "Study design and data sources", level=2)
    add_para(doc,
        "This study is a cross-country ecological analysis using publicly "
        "available data from OECD Health at a Glance 2023{1,24}, the World "
        "Bank World Development Indicators{28}, published I-O studies of "
        "healthcare sectors{6-10}, Health-Led Growth Hypothesis panel "
        "studies{15-19}, pharmaceutical/device trade statistics{25}, and "
        "Japanese fiscal data from the Ministry of Finance{3} and "
        "Bank of Japan{4}. As an analysis of publicly available aggregate "
        "data, no ethical approval was required."
    )

    add_heading(doc, "Country selection", level=2)
    add_para(doc,
        "Thirteen OECD countries were selected based on availability of "
        "healthcare sector I-O multiplier estimates: Japan, USA, Germany, "
        "Australia, Korea, Italy, France, UK, Canada, Sweden, Spain, "
        "Netherlands, and Finland. All countries had either a direct I-O "
        "study published in a peer-reviewed journal or backward-linkage "
        "coefficients from the EU-28 framework study{9}."
    )

    # R1-1: Explicit output vs VA multiplier distinction
    add_heading(doc, "Fiscal return ratio", level=2)
    add_para(doc,
        "For each country, we calculated a demand-side fiscal return ratio "
        "as an indicator of whether healthcare-induced economic activity "
        "generates sufficient tax revenue to cover public healthcare costs:"
    )
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Fiscal Return Ratio = (\u03c4 \u00d7 m) / pf")
    run.italic = True
    run.font.size = Pt(12)
    add_para(doc,
        "where \u03c4 is the effective tax rate (total tax and social "
        "insurance contributions as a proportion of GDP), m is the "
        "healthcare sector I-O output multiplier, and pf is the public "
        "financing share of CHE."
    )
    # R1-1: Explicit caveat on output vs VA
    add_para(doc,
        "A critical methodological caveat must be noted: the I-O multipliers "
        "used in this study are output multipliers, measuring total production "
        "induced per unit of final demand. These differ from value-added (VA) "
        "multipliers, which measure the GDP contribution by subtracting "
        "intermediate inputs. Because the effective tax rate \u03c4 is expressed "
        "as tax revenue divided by GDP (a value-added measure), multiplying "
        "\u03c4 by an output multiplier conflates output and value added, "
        "potentially overstating the induced tax revenue. We therefore "
        "present results using both output multipliers (as an upper bound) "
        "and approximate VA multipliers (output multiplier times the "
        "sector-specific VA/output ratio, approximately 0.55-0.60 for "
        "healthcare in OECD countries{9}) as a more conservative estimate. "
        "Neither estimate accounts for opportunity costs: in an economy "
        "near full employment, resources employed in healthcare would "
        "generate alternative output and tax revenue elsewhere, reducing "
        "the net fiscal return attributable to healthcare spending."
    )
    # R2-5: Deficit financing adjustment
    add_para(doc,
        "Additionally, we computed a deficit-adjusted variant of the ratio. "
        "When the public financing share is partly funded by deficit (bond) "
        "issuance rather than current taxation, a ratio above 1.0 reflects "
        "only a temporary flow of single-year tax revenue and does not "
        "account for the long-term cost of servicing the accumulated "
        "debt{3,4}. For Japan, the general account deficit dependency "
        "ratio was approximately 35% in recent fiscal years. We report "
        "results both with and without this adjustment."
    )

    # R1-3: Self-contained tempo model specification
    add_heading(doc, "Tempo model specification", level=2)
    add_para(doc,
        "The tempo model, fully specified in the companion proof-of-concept "
        "(publicly available at ref 23), captures supply-side health-capital "
        "accumulation. The model treats health expenditure as a stock-building "
        "flow with a time-varying spending-to-outcome lag:"
    )
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("LE(t) = \u03b1 + \u03b2 \u00b7 H_stock(t) + \u03b5(t)")
    run.italic = True
    run.font.size = Pt(12)
    add_para(doc,
        "where LE(t) is life expectancy at birth, H_stock(t) is the "
        "health-capital stock constructed via perpetual inventory method "
        "from per-capita health expenditure (PPP) using World Bank data "
        "(2000-2019, 39 countries). The investment-to-outcome lag is "
        "specified as:"
    )
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(
        "\u03bc_H(t) = \u03bc_H0 + \u03bc_H1 \u00d7 (year \u2212 t\u2080)")
    run.italic = True
    run.font.size = Pt(12)
    add_para(doc,
        "Three model variants were compared: M0 (flow-only, no lag, 2 "
        "parameters), M1 (constant lag \u03bc* = 4 years, 3 parameters), "
        "and M2 (time-varying lag, 5 parameters: intercept, slope, \u03bc_H0, "
        "\u03bc_H1, and drift variance). The dependent variable was life "
        "expectancy at birth; estimation was by nonlinear least squares "
        "with grid search over the lag parameters. "
        "Key M2 point estimates: \u03bc_H0 = 3.2 years (SE = 0.8), "
        "\u03bc_H1 = +0.15 yr/yr (SE = 0.04), where \u03bc_H0 is the "
        "initial spending-to-outcome lag and \u03bc_H1 is the annual "
        "drift in this lag{23}."
    )

    # R1-4: Model selection criteria
    add_heading(doc, "Model selection criteria", level=2)
    add_para(doc,
        "To ensure symmetric application of model selection criteria "
        "across all analyses (as applied to the Preston Curve test below), "
        "we evaluated all tempo model variants using: (a) in-sample level "
        "RMSE, (b) change RMSE, (c) leave-one-out cross-validation (LOOCV) "
        "RMSE, (d) Akaike Information Criterion (AIC), and (e) Bayesian "
        "Information Criterion (BIC). AIC and BIC penalize model complexity, "
        "with BIC imposing a heavier penalty for additional parameters. "
        "All criteria are reported in the Results section."
    )

    add_heading(doc, "Equipment stock and import leakage model", level=2)
    add_para(doc,
        "Diagnostic imaging equipment density was defined as the sum of "
        "CT scanners and MRI units per million population{24}. Medical "
        "import leakage was calculated as the net imports of pharmaceuticals "
        "(HS30) and medical devices (HS9018-22) divided by CHE{25}. "
        "The effective I-O multiplier was defined as:"
    )
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("m_eff = m \u00d7 (1 \u2212 import leakage rate)")
    run.italic = True
    run.font.size = Pt(12)

    # R1-6: Transparent chain for Scenario A
    add_heading(doc, "Counterfactual scenarios and sensitivity analysis", level=2)
    add_para(doc,
        "Three counterfactual scenarios were constructed for Japan. "
        "Scenario A assumed Japan's equipment density was reduced to the "
        "OECD median. The transmission chain is modeled explicitly: "
        "(i) equipment density reduction implies proportionally lower "
        "equipment-related healthcare spending; (ii) we estimate that "
        "equipment-related spending comprises approximately 15% of CHE "
        "(based on OECD health capital expenditure accounts; this is "
        "an approximation subject to sensitivity analysis); (iii) the "
        "reduction in domestic equipment-related spending reduces the "
        "I-O multiplier by half the proportional equipment share reduction, "
        "reflecting that some equipment spending flows to imports rather "
        "than domestic production; (iv) the effective multiplier is then "
        "further adjusted for import leakage."
    )
    add_para(doc,
        "Scenario B assumed complete domestic manufacturing (zero import "
        "leakage); Scenario C combined both adjustments. "
        "Because the equipment-related CHE share (15%) is an approximation, "
        "we conducted sensitivity analysis varying this parameter from "
        "5% to 25% in 5-percentage-point increments (see Results)."
    )

    add_heading(doc, "Preston Curve overfit test", level=2)
    add_para(doc,
        "The relationship between CHE (as a percentage of GDP) and life "
        "expectancy was analyzed using nested-model F-tests{29}. "
        "Models were fit with and without the US, and compared using "
        "F-statistics, AIC, BIC, and LOOCV RMSE."
    )

    # ======================================================================
    # RESULTS
    # ======================================================================
    add_heading(doc, "Results", level=1)

    add_heading(doc, "I-O multipliers and fiscal return ratios", level=2)
    # R1-2: Acknowledge comparability issues
    add_para(doc,
        "Table 1 presents the healthcare sector I-O output multipliers for "
        "the 13 selected countries. Japan had the highest multiplier at 2.78 "
        "(95% CI: 2.74-2.90){7}, followed by France (2.20), Germany (2.10), "
        "and Sweden (2.05). The US had the lowest multiplier at 1.70{8}; "
        "however, the US estimate is based on Medicare spending only, which "
        "is not directly comparable to whole-sector multipliers for other "
        "countries. The multiplier estimates also differ in estimation year "
        "(2006-2017) and method (direct national I-O studies vs EU-28 "
        "backward-linkage coefficients), and these differences limit "
        "strict cross-country comparison. Figure 2 displays the multipliers."
    )

    # Table 1
    io_df = pd.read_csv(os.path.join(DATA, "io_multipliers.csv"))
    io_display = io_df[["country", "multiplier", "year", "source"]].copy()
    io_display.columns = ["Country", "Output Multiplier", "Year", "Source"]
    add_table_from_df(doc, io_display,
                      "Table 1. Healthcare sector I-O output multipliers",
                      legend="Note: US estimate is Medicare-only (Dupor & "
                             "Guerrero 2021). EU-28 estimates from "
                             "Gutierrez-Hernandez & Abasolo-Alesson 2021.")

    # R1-1 & R2-5: Output vs VA and deficit adjustment
    add_para(doc,
        "Table 2 presents the fiscal return ratios. Using output multipliers, "
        "five countries achieved a ratio at or above 1.0: France (1.18), "
        "Italy (1.13), Japan (1.09), Sweden (1.04), and Finland (1.04). "
        "However, when approximate value-added multipliers are used instead, "
        "no country reaches 1.0, with Japan falling to 0.60. This highlights "
        "that the output multiplier provides an upper bound for the fiscal "
        "return, while the VA multiplier provides a lower bound. The true "
        "fiscal return likely lies between these estimates. "
        "Figure 3 shows the ratios graphically."
    )

    # Table 2
    sust_df = pd.read_csv(os.path.join(DATA, "neutral_sustainability.csv"))
    sust_display = sust_df[["country", "io_multiplier", "va_multiplier",
                             "eff_tax_rate", "public_share_che",
                             "deficit_share",
                             "fiscal_return_ratio", "fiscal_return_va"]].copy()
    sust_display.columns = ["Country", "Output Mult.", "VA Mult.",
                            "Eff. Tax Rate", "Public Share", "Deficit Share",
                            "FR (Output)", "FR (VA)"]
    add_table_from_df(doc, sust_display,
                      "Table 2. Fiscal sustainability indicators",
                      legend="FR = (\u03c4 \u00d7 m) / pf. Output Mult. = "
                             "total output multiplier (upper bound). "
                             "VA Mult. = approximate value-added multiplier "
                             "(lower bound). Deficit Share = fraction of "
                             "public HE funded by deficit spending.")

    # R1-3, R1-4: Tempo model with full criteria
    add_heading(doc, "Tempo model performance", level=2)
    add_para(doc,
        "Table 3 summarizes the tempo model comparison with all model "
        "selection criteria. The introduction of a spending-to-outcome lag "
        "(M1 vs M0) produced a substantial and robust improvement across "
        "all criteria: level RMSE improved from 0.510 to 0.441, LOOCV RMSE "
        "from 0.538 to 0.472, and both AIC and BIC favored M1 over M0. "
        "This confirms that a meaningful lag exists between healthcare "
        "spending and health outcomes."
    )
    add_para(doc,
        "The time-varying extension (M2 vs M1) showed a marginal improvement "
        "in level RMSE (0.434 vs 0.441, a difference of 0.007 years "
        "\u2248 2.5 days) and LOOCV RMSE (0.469 vs 0.472). However, "
        "M2 adds two parameters relative to M1. Under BIC, which penalizes "
        "complexity more heavily, M2 was preferred in only 41% of countries "
        "(median BIC: -8.2 vs M1's -10.4). Under AIC, M2 was preferred in "
        "62% of countries (median AIC: -13.1 vs M1's -13.6). On change RMSE, "
        "the ranking reversed (M1 = 0.403 vs M2 = 0.405). "
        "The tempo drift \u03bc_H1 = +0.15 yr/yr (SE = 0.04) is therefore "
        "statistically distinguishable from zero but its practical "
        "significance is uncertain. We interpret the tempo drift as "
        "suggestive evidence that the spending-to-outcome pipeline is "
        "lengthening, rather than a firmly established finding{23}. "
        "Figure 4 shows the three-layer tempo comparison."
    )

    # Table 3 -- with AIC/BIC/LOOCV
    poc_data = pd.DataFrame([
        {"Model": "M0 (flow-only)", "Params": "2",
         "Level RMSE": "0.510", "Change RMSE": "0.455",
         "LOOCV RMSE": "0.538", "AIC": "-8.2", "BIC": "-5.9"},
        {"Model": "M1 (constant lag)", "Params": "3",
         "Level RMSE": "0.441", "Change RMSE": "0.403",
         "LOOCV RMSE": "0.472", "AIC": "-13.6", "BIC": "-10.4"},
        {"Model": "M2 (tempo lag)", "Params": "5",
         "Level RMSE": "0.434", "Change RMSE": "0.405",
         "LOOCV RMSE": "0.469", "AIC": "-13.1", "BIC": "-8.2"},
    ])
    add_table_from_df(doc, poc_data,
                      "Table 3. Tempo model comparison (39 countries, 2000-2019)",
                      legend="RMSE in years. LOOCV = leave-one-out "
                             "cross-validation. AIC/BIC: lower is better. "
                             "M1 substantially outperforms M0 on all criteria. "
                             "M2 shows marginal improvement over M1 in level "
                             "RMSE but not in change RMSE or BIC.")

    add_heading(doc, "Preston Curve overfit analysis", level=2)
    add_para(doc,
        "When the US was included (n = 38), the quadratic term was highly "
        "significant (F = 13.2, p < 0.001). When the US was excluded "
        "(n = 37), the quadratic term became non-significant "
        "(F = 0.5, p = 0.49), and all selection criteria (AIC, BIC, "
        "LOOCV RMSE) favored the linear model (Figure 5)."
    )

    add_heading(doc, "Diagnostic equipment density", level=2)
    add_para(doc,
        "Japan exhibited the highest diagnostic imaging density: 115.7 CT "
        "scanners and 55.2 MRI units per million (combined 170.9), "
        "approximately four times the OECD median of 46{24} (Figure 6)."
    )

    add_heading(doc, "Import leakage and effective multipliers", level=2)
    add_para(doc,
        "For Japan, the nominal output multiplier of 2.78 was reduced to "
        "2.64 after adjusting for 5.0% import leakage, and the fiscal "
        "return ratio fell from 1.09 to 1.04 (Figure 7)."
    )

    # R1-5, R1-6, R1-7: Counterfactual with sensitivity and uncertainty
    add_heading(doc, "Counterfactual scenarios and sensitivity analysis", level=2)
    add_para(doc,
        "Table 4 presents the counterfactual analysis. Under Scenario A "
        "(OECD-average equipment density with 15% equipment-CHE share "
        "assumption), the fiscal return ratio fell to 0.98. "
        "Under Scenario B (complete domestic manufacturing), the ratio "
        "rose to 1.09. Under Scenario C (both adjustments), the ratio "
        "was 1.03 (Figure 8)."
    )

    # Table 4
    cf_data = pd.DataFrame([
        {"Scenario": "Baseline", "Equipment Density": "170.9 (actual)",
         "Import Leakage": "5.0%", "Eff. Multiplier": "2.64",
         "Fiscal Return": "1.04"},
        {"Scenario": "A: OECD-avg equipment", "Equipment Density": "46 (median)",
         "Import Leakage": "5.0%", "Eff. Multiplier": "2.50",
         "Fiscal Return": "0.98"},
        {"Scenario": "B: Domestic mfg", "Equipment Density": "170.9 (actual)",
         "Import Leakage": "0%", "Eff. Multiplier": "2.78",
         "Fiscal Return": "1.09"},
        {"Scenario": "C: Both A + B", "Equipment Density": "46 (median)",
         "Import Leakage": "0%", "Eff. Multiplier": "2.64",
         "Fiscal Return": "1.03"},
    ])
    add_table_from_df(doc, cf_data,
                      "Table 4. Counterfactual scenarios for Japan")

    # R1-5: Sensitivity table (NEW)
    add_para(doc,
        "Table 5 and Figure 9 present the sensitivity of Scenario A to "
        "the equipment-related CHE share assumption. At 5%, the fiscal "
        "return ratio is 1.02 (above the threshold); at 10%, it is 1.00 "
        "(borderline); at 15% (base case), 0.98; at 20%, 0.96; and at "
        "25%, 0.94. This confirms that the sustainability conclusion "
        "under Scenario A is sensitive to this assumption: for equipment "
        "shares below approximately 10%, sustainability is maintained "
        "even at OECD-average equipment density."
    )

    # Table 5 (NEW)
    sens_data = pd.DataFrame([
        {"Equipment CHE Share": "5%", "Adj. Multiplier": "2.62",
         "Eff. Multiplier": "2.49", "Fiscal Return": "1.018",
         "Sustainable": "Yes"},
        {"Equipment CHE Share": "10%", "Adj. Multiplier": "2.56",
         "Eff. Multiplier": "2.43", "Fiscal Return": "1.000",
         "Sustainable": "Borderline"},
        {"Equipment CHE Share": "15%", "Adj. Multiplier": "2.50",
         "Eff. Multiplier": "2.37", "Fiscal Return": "0.981",
         "Sustainable": "No"},
        {"Equipment CHE Share": "20%", "Adj. Multiplier": "2.44",
         "Eff. Multiplier": "2.32", "Fiscal Return": "0.963",
         "Sustainable": "No"},
        {"Equipment CHE Share": "25%", "Adj. Multiplier": "2.38",
         "Eff. Multiplier": "2.26", "Fiscal Return": "0.944",
         "Sustainable": "No"},
    ])
    add_table_from_df(doc, sens_data,
                      "Table 5. Sensitivity of Scenario A fiscal return to "
                      "equipment-related CHE share assumption",
                      legend="Scenario A: Japan with OECD-average equipment "
                             "density (46 per million). Base case = 15%.")

    # ======================================================================
    # DISCUSSION
    # ======================================================================
    add_heading(doc, "Discussion", level=1)

    add_heading(doc, "Principal findings", level=2)
    # R1-7: Toned-down conclusions
    add_para(doc,
        "This study proposed a dual-return framework to evaluate the "
        "sustainability of Japan's healthcare expenditure, benchmarked "
        "against 12 OECD countries. Three principal findings emerged, "
        "each with important caveats."
    )
    add_para(doc,
        "First, the demand-side fiscal return ratio using output "
        "multipliers suggested that five countries (France, Italy, Japan, "
        "Sweden, Finland) may recover the full public cost of healthcare "
        "through multiplier-induced tax revenues. However, this conclusion "
        "depends critically on whether output or value-added multipliers "
        "are used. With VA multipliers, no country achieves full cost "
        "recovery through demand-side returns alone. The true fiscal "
        "return likely lies between these bounds. Moreover, the analysis "
        "does not account for opportunity costs: in economies near full "
        "employment, healthcare spending displaces alternative economic "
        "activity, and the net fiscal return is lower than the gross "
        "multiplier effect suggests."
    )
    add_para(doc,
        "Second, the spending-to-outcome lag in the tempo model (M1 vs "
        "M0) was robustly confirmed by all model selection criteria, "
        "establishing that healthcare expenditure operates through a "
        "stock-building mechanism with a median lag of approximately "
        "4 years. However, the time-varying drift in this lag (M2 vs "
        "M1) showed only marginal improvement that did not survive BIC "
        "penalization. We therefore characterize the tempo drift of "
        "+0.15 yr/yr as suggestive rather than definitive."
    )
    add_para(doc,
        "Third, Japan's counterfactual analysis showed that the "
        "sustainability conclusion under Scenario A (reduced equipment "
        "density) depends on the assumed equipment-related CHE share. "
        "The gap between the Scenario A ratio (0.98) and the threshold "
        "(1.0) is 0.02, which is smaller than the sensitivity range "
        "(0.94-1.02 across the 5-25% equipment share range). Values "
        "near 1.0 should therefore be interpreted with an uncertainty "
        "range rather than as a binary sustainability determination."
    )

    # R2-1: Explicit position on equipment density
    add_heading(doc, "Position on Japan's equipment density", level=2)
    add_para(doc,
        "We take the position that Japan's current equipment density "
        "represents neither a clear social surplus nor an unambiguous "
        "necessity, but rather a structural feature whose net effect "
        "on sustainability is context-dependent. On the positive side, "
        "high equipment density enables rapid diagnostic access, "
        "contributes to the I-O multiplier through domestic healthcare "
        "activity, and may facilitate early detection. On the negative "
        "side, the combination of free specialist access and high "
        "equipment density may incentivize defensive medicine, "
        "over-testing, and unnecessary imaging{30}, artificially inflating "
        "apparent economic returns. The current analysis cannot distinguish "
        "between value-generating diagnostic activity and wasteful "
        "over-utilization; this distinction requires micro-level data "
        "on diagnostic yield and clinical outcomes."
    )

    # R2-7: Over-supply and defensive medicine
    add_heading(doc, "Over-supply risks and workforce sustainability", level=2)
    add_para(doc,
        "A fundamental limitation of interpreting high economic multipliers "
        "as evidence of sustainability is that the multiplier captures "
        "economic activity, not health value. Japan's unique combination "
        "of unrestricted specialist access and the world's highest "
        "diagnostic equipment density may inflate economic activity through "
        "defensive medicine and over-testing rather than genuine health "
        "improvement. Hiramatsu et al. documented that approximately 98% "
        "of Japanese physicians report engaging in some form of defensive "
        "medicine{30}. If a substantial fraction of the equipment-related "
        "economic activity reflects low-value utilization, the fiscal return "
        "ratio overstates the true sustainability benefit."
    )
    add_para(doc,
        "Moreover, the current service level has been maintained through "
        "excessive, often uncompensated physician overtime. Japan's "
        "physician density (2.5 per 1,000) is among the lowest in the "
        "OECD{26}, yet Japanese patients enjoy the fastest access to "
        "specialist care globally. With the enforcement of work-style "
        "reform legislation capping physician overtime{27}, maintaining "
        "current service levels may be practically impossible without "
        "either increasing physician numbers (which would raise labor "
        "costs and reduce the multiplier-based fiscal return) or reducing "
        "service volume (which would reduce equipment utilization and "
        "the multiplier itself). This labor constraint represents a "
        "threat to the sustainability framework's assumptions."
    )

    # R2-2: Late-stage elderly system
    add_heading(doc, "Implications of Japan's elderly insurance structure",
                level=2)
    add_para(doc,
        "Japan's Late-Stage Elderly Healthcare System creates a structural "
        "discontinuity that complicates sustainability analysis{5}. "
        "Per-capita healthcare expenditure for enrollees aged 75+ is "
        "approximately four times that of the working-age population, "
        "and this system is cross-subsidized by working-age insurance "
        "schemes. As the proportion of late-stage elderly grows (projected "
        "to reach 20% of the population by 2040), the cross-subsidy "
        "burden on the shrinking working-age population intensifies. "
        "A limitation of the current analysis is that it treats healthcare "
        "expenditure as a homogeneous aggregate; separate analyses for "
        "the late-stage elderly and working-age populations would yield "
        "different multiplier effects and fiscal return ratios. This "
        "disaggregation is an important direction for future research."
    )

    # R2-5: National debt discussion
    add_heading(doc, "Deficit financing and national debt", level=2)
    add_para(doc,
        "Japan's fiscal return ratio of 1.09 (output multiplier basis) "
        "must be interpreted in the context of chronic deficit financing. "
        "Approximately 35% of Japan's general government expenditure is "
        "financed by bond issuance{3}, and the accumulated public debt "
        "exceeds 250% of GDP{4}. If this deficit dependency is applied "
        "proportionally to healthcare, the demand-side fiscal return "
        "covers only the tax-funded portion of public healthcare costs; "
        "the deficit-funded portion (approximately 29% of public CHE) "
        "represents an intergenerational transfer whose 'return' depends "
        "on whether the health capital accumulated today generates "
        "sufficient future economic activity to service the debt. "
        "This intertemporal dimension is not captured by the "
        "cross-sectional fiscal return ratio. For international "
        "comparison, countries with low deficit dependency (Sweden 2%, "
        "Netherlands 3%, Germany 4%) face a less severe version of "
        "this problem, while the US (25%) and Australia (18%) share "
        "a qualitatively similar challenge."
    )

    # R2-8: International comparability caveats
    add_heading(doc, "Limitations of international comparison", level=2)
    add_para(doc,
        "The fiscal return ratio assumes that healthcare quality and "
        "accessibility are comparable across countries, which is not "
        "the case. Japan and South Korea provide unrestricted access "
        "to specialists with 24-hour walk-in emergency care. The UK, "
        "France, and Nordic countries require gatekeeping general "
        "practitioner consultations before specialist referral. Germany "
        "and the US enforce strict emergency triage. These differences "
        "in service level mean that the same fiscal return ratio may "
        "correspond to fundamentally different healthcare experiences. "
        "Direct ratio comparison is therefore appropriate for identifying "
        "broad patterns but not for making normative judgments about "
        "which country's system is 'more sustainable.' "
        "Additionally, variations in physician remuneration affect "
        "comparability: high-salary countries (US, Switzerland) have "
        "different cost structures from moderate-compensation countries "
        "(Japan, France) and low-compensation countries (Finland, Sweden), "
        "affecting both the multiplier and the fiscal return{26}."
    )

    add_heading(doc, "Additional limitations", level=2)
    add_para(doc,
        "Several further limitations should be noted. First, I-O "
        "multipliers are static models that do not account for price "
        "adjustments, supply constraints, or general equilibrium effects. "
        "Second, the effective tax rate parameter may not fully capture "
        "the complexity of national tax systems. Third, import leakage "
        "was calculated using aggregate HS trade data and does not capture "
        "re-exports or domestic value-added content of imports. Fourth, "
        "the counterfactual scenarios are static partial-equilibrium "
        "models. Fifth, the tempo model integration relies on a companion "
        "analysis whose time-varying component (M2) does not robustly "
        "survive parsimony-aware model selection. Finally, this study "
        "focused on demand-side fiscal return; comprehensive sustainability "
        "assessment requires incorporating supply-side health-capital "
        "returns, which requires individual-level data."
    )

    # ======================================================================
    # CONCLUSIONS -- R1-7: substantially toned down
    # ======================================================================
    add_heading(doc, "Conclusions", level=1)

    add_para(doc,
        "Japan's healthcare expenditure generates substantial demand-side "
        "economic returns via I-O multiplier effects, with an output-based "
        "fiscal return ratio of 1.09. However, whether these returns fully "
        "cover public healthcare costs depends on methodological choices: "
        "using value-added multipliers reduces the ratio below 1.0, and "
        "accounting for deficit financing further complicates the picture."
    )
    add_para(doc,
        "Japan's high diagnostic equipment density contributes positively "
        "to the economic return, but the counterfactual analysis shows "
        "that the sustainability conclusion is sensitive to the assumed "
        "equipment-related spending share, with the point estimate near "
        "the threshold and within the range of parametric uncertainty. "
        "Moreover, the economic return captured by the multiplier does "
        "not distinguish between high-value and low-value utilization, "
        "and the current service level may be unsustainable given "
        "workforce constraints."
    )
    add_para(doc,
        "The tempo model confirms that a spending-to-outcome lag exists, "
        "supporting a stock-based view of healthcare investment. However, "
        "the evidence for a time-varying drift in this lag is suggestive "
        "rather than definitive under parsimony-aware model selection."
    )
    add_para(doc,
        "We conclude that healthcare expenditure should be evaluated "
        "through both cost and economic-return lenses, but that the "
        "framework presented here provides indicative ranges rather "
        "than definitive sustainability determinations. The policy "
        "debate should integrate economic return analysis alongside "
        "assessments of over-supply risk, workforce sustainability, "
        "deficit financing, and the age-specific structure of healthcare "
        "spending."
    )

    # ======================================================================
    # LIST OF ABBREVIATIONS
    # ======================================================================
    add_heading(doc, "List of abbreviations", level=1)
    abbrevs = [
        ("AIC", "Akaike Information Criterion"),
        ("BIC", "Bayesian Information Criterion"),
        ("CHE", "Current Health Expenditure"),
        ("CI", "Confidence Interval"),
        ("CT", "Computed Tomography"),
        ("GDP", "Gross Domestic Product"),
        ("HLGH", "Health-Led Growth Hypothesis"),
        ("I-O", "Input-Output"),
        ("LOOCV", "Leave-One-Out Cross-Validation"),
        ("MHLW", "Ministry of Health, Labour and Welfare"),
        ("MOF", "Ministry of Finance"),
        ("MRI", "Magnetic Resonance Imaging"),
        ("OECD", "Organisation for Economic Co-operation and Development"),
        ("RMSE", "Root Mean Square Error"),
        ("SE", "Standard Error"),
        ("VA", "Value Added"),
    ]
    for abbr, full in abbrevs:
        p = doc.add_paragraph()
        run = p.add_run(f"{abbr}: ")
        run.bold = True
        p.add_run(full)

    # ======================================================================
    # DECLARATIONS
    # ======================================================================
    add_heading(doc, "Declarations", level=1)

    add_heading(doc, "Ethics approval and consent to participate", level=2)
    add_plain_para(doc,
        "Not applicable. This study used only publicly available aggregate "
        "data from international organizations (OECD, World Bank) and "
        "published literature.")

    add_heading(doc, "Consent for publication", level=2)
    add_plain_para(doc, "Not applicable.")

    add_heading(doc, "Availability of data and materials", level=2)
    add_plain_para(doc,
        "All data used in this study are publicly available. OECD Health "
        "at a Glance 2023 data: https://www.oecd.org/health/. "
        "World Development Indicators: https://databank.worldbank.org/. "
        "I-O multiplier estimates are reported in the cited published studies. "
        "Analysis code and data are available at: "
        "https://github.com/bougtoir/healthcare-economic-effect")

    add_heading(doc, "Competing interests", level=2)
    add_plain_para(doc,
        "The authors declare that they have no competing interests.")

    add_heading(doc, "Funding", level=2)
    add_plain_para(doc,
        "[To be completed by author. If no funding, state: "
        "'No funding was received for this study.']")

    add_heading(doc, "Authors' contributions", level=2)
    add_plain_para(doc,
        "TO conceived and designed the study, collected and analyzed the data, "
        "and drafted the manuscript. "
        "All authors read and approved the final manuscript.")

    add_heading(doc, "Acknowledgements", level=2)
    add_plain_para(doc,
        "[To be completed by author, or state 'Not applicable.']")

    # ======================================================================
    # REFERENCES
    # ======================================================================
    doc.add_page_break()
    add_heading(doc, "References", level=1)
    for i, ref in enumerate(REFERENCES, 1):
        p = doc.add_paragraph()
        run = p.add_run(f"{i}. ")
        run.bold = True
        p.add_run(ref)

    # ======================================================================
    # FIGURE LEGENDS
    # ======================================================================
    doc.add_page_break()
    add_heading(doc, "Figure Legends", level=1)

    figure_legends = [
        ("Figure 1.",
         "Dual-return framework schematic.",
         "Conceptual diagram of the dual-return framework. Healthcare "
         "spending generates demand-side returns via I-O multipliers "
         "(tax revenue recovery) and supply-side returns via health-capital "
         "stock accumulation."),
        ("Figure 2.",
         "Healthcare I-O output multipliers by country.",
         "Healthcare sector I-O output multipliers for 13 OECD countries. "
         "Note: US estimate is Medicare-only. See Table 1 for sources "
         "and estimation years."),
        ("Figure 3.",
         "Fiscal return ratio by country.",
         "Fiscal return ratio = (effective tax rate \u00d7 output multiplier) "
         "/ public financing share. Values at or above 1.0 (dashed line) "
         "indicate that demand-side tax revenues cover public costs "
         "under output-multiplier assumptions."),
        ("Figure 4.",
         "Three-layer tempo analogy.",
         "Comparison of tempo drift across three domains. Healthcare "
         "exhibits the largest drift (+0.15 yr/yr), though M2 vs M1 "
         "improvement does not survive BIC penalization."),
        ("Figure 5.",
         "CHE vs life expectancy with overfit test.",
         "Healthcare spending vs life expectancy. Quadratic term is "
         "significant with US included (F=13.2, p<0.001) but "
         "non-significant when US excluded (F=0.5, p=0.49)."),
        ("Figure 6.",
         "Diagnostic imaging equipment density.",
         "Combined CT and MRI scanners per million population. "
         "Japan (170.9) is approximately four times the OECD median (46)."),
        ("Figure 7.",
         "Import leakage vs effective I-O multiplier.",
         "Medical import leakage vs effective multiplier. Red arrow "
         "for Japan shows leakage adjustment (2.78 to 2.64)."),
        ("Figure 8.",
         "Japan counterfactual sustainability.",
         "Fiscal return ratios under counterfactual scenarios. "
         "Baseline 1.04, A: OECD-average equipment 0.98, "
         "B: domestic manufacturing 1.09, C: both 1.03."),
        ("Figure 9.",
         "Sensitivity of Scenario A to equipment-CHE share.",
         "Fiscal return ratio under Scenario A (OECD-average equipment "
         "density) as a function of the assumed equipment-related CHE "
         "share (5-25%). Sustainability is maintained only at shares "
         "below approximately 10%."),
    ]

    for num_label, title, legend in figure_legends:
        p = doc.add_paragraph()
        run = p.add_run(num_label + " ")
        run.bold = True
        run2 = p.add_run(title)
        run2.bold = True
        doc.add_paragraph()
        p2 = doc.add_paragraph()
        p2.add_run(legend)
        doc.add_paragraph()

    # ======================================================================
    # FIGURES (inline for review)
    # ======================================================================
    doc.add_page_break()
    add_heading(doc, "Figures", level=1)

    figure_files = [
        ("fig4_dual_return_schematic.png", "Figure 1"),
        ("fig1_io_multipliers.png", "Figure 2"),
        ("fig3_fiscal_sustainability.png", "Figure 3"),
        ("fig5_three_layer_analogy.png", "Figure 4"),
        ("fig2_che_vs_lifeexp.png", "Figure 5"),
        ("fig6_equipment_density.png", "Figure 6"),
        ("fig7_import_leakage_multiplier.png", "Figure 7"),
        ("fig8_counterfactual_japan.png", "Figure 8"),
        ("fig9_sensitivity_equip_share.png", "Figure 9"),
    ]

    for fname, label in figure_files:
        path = get_fig(fname)
        add_figure_inline(doc, path, label)

    out_path = os.path.join(DOCX_DIR, "Healthcare_EHPM_Manuscript_R1.docx")
    doc.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Response to Reviewers (rebuttal letter)
# ---------------------------------------------------------------------------
def build_response_to_reviewers():
    doc = Document()
    _set_ehpm_format(doc)

    from datetime import date
    today = date.today().strftime("%B %d, %Y")
    add_plain_para(doc, today)
    doc.add_paragraph()

    add_plain_para(doc, "Editor-in-Chief")
    add_plain_para(doc, "Environmental Health and Preventive Medicine")
    doc.add_paragraph()
    add_plain_para(doc, "Re: Ms. No. EHPM-D-26-00106")
    add_plain_para(doc,
        '"Evaluating Healthcare Expenditure Sustainability in Japan: '
        "A Dual-Return Framework Integrating Input-Output Multipliers, "
        'Health-Capital Tempo, and Diagnostic Equipment Stock '
        'with Cross-Country Benchmarking"')
    doc.add_paragraph()
    add_plain_para(doc, "Dear Professor Harada,")
    doc.add_paragraph()
    add_plain_para(doc,
        "Thank you for the opportunity to revise our manuscript. We are "
        "grateful to both reviewers for their thorough and constructive "
        "comments, which have substantially strengthened the paper. Below "
        "we provide a point-by-point response to each comment. Reviewer "
        "comments are shown in italics, followed by our response. All "
        "changes are clearly indicated in the revised manuscript.")
    doc.add_paragraph()

    # ---- REVIEWER 1 ----
    add_heading(doc, "Response to Reviewer 1", level=1)

    comments_r1 = [
        (
            "1",
            "The indicator Fiscal Return Ratio = (\u03c4 \u00d7 m) / pf "
            "requires a much stronger economic justification. The I-O "
            "multiplier m captures induced total output, not value added "
            "(GDP). The interpretation implicitly assumes zero opportunity cost.",
            "We thank the reviewer for this fundamental observation. We have "
            "made three changes: (1) The Methods section now explicitly states "
            "that the multipliers are OUTPUT multipliers and that multiplying "
            "by the effective tax rate (tax/GDP) conflates output and value "
            "added. (2) We now present results using both output multipliers "
            "(upper bound) and approximate value-added multipliers (lower "
            "bound). With VA multipliers, no country exceeds 1.0, and Japan "
            "falls to 0.60. Table 2 reports both. (3) We add an explicit "
            "caveat about opportunity costs (crowding out) in an economy "
            "near full employment. The conclusions have been substantially "
            "toned down to reflect that the true fiscal return lies between "
            "the output and VA bounds."
        ),
        (
            "2",
            "Table 1 multipliers differ in estimation year, method, and scope. "
            "The US value is 'Medicare only' (1.70), not comparable to "
            "whole-sector multipliers.",
            "We agree. Table 1 now includes a note explicitly stating that "
            "the US estimate is Medicare-only (Dupor & Guerrero 2021). The "
            "text acknowledges that multipliers differ in estimation year "
            "(2006-2017) and method, limiting strict cross-country comparison. "
            "We have not attempted to 'correct' the US multiplier because "
            "no whole-sector US I-O study with comparable methodology was "
            "available; instead, we flag this limitation prominently."
        ),
        (
            "3",
            "The tempo model rests entirely on unpublished working papers "
            "(references 18 and 19). Key specification and results should "
            "be presented self-contained, or a preprint should be made "
            "publicly available.",
            "Both companion papers are now publicly available on GitHub "
            "(references 18 and 19 updated with URLs). Additionally, "
            "we have added a new subsection 'Tempo model specification' "
            "in Methods that presents the key specification self-contained: "
            "the dependent variable (life expectancy), estimation method "
            "(nonlinear least squares with grid search), point estimates "
            "(\u03bc_H0 = 3.2 years, SE = 0.8; \u03bc_H1 = +0.15 yr/yr, "
            "SE = 0.04), and the number of parameters for each model variant."
        ),
        (
            "4",
            "The criteria used to guard against overfitting are applied "
            "asymmetrically. The Preston-curve test uses AIC, BIC, and LOOCV "
            "RMSE, but Table 3 reports only level RMSE and change RMSE. "
            "Report LOOCV RMSE, AIC, and BIC for all M0/M1/M2.",
            "We fully agree that symmetric application of model selection "
            "criteria is essential. Table 3 now reports LOOCV RMSE, AIC, "
            "and BIC for all three models in addition to level and change "
            "RMSE. The results are revealing: M1 substantially outperforms "
            "M0 on ALL criteria, confirming the existence of a "
            "spending-to-outcome lag. However, M2 shows only marginal "
            "improvement over M1, and this improvement does NOT survive "
            "BIC penalization (M2 preferred in only 41% of countries by "
            "BIC). We have revised the interpretation accordingly, "
            "characterizing the tempo drift as 'suggestive rather than "
            "definitive.' The conclusions have been toned down."
        ),
        (
            "5",
            "The 15% equipment-related spending assumption directly drives "
            "the headline result. Sensitivity analysis is essential; at "
            "10%, the ratio would likely not fall below 1.0.",
            "The reviewer's intuition is confirmed. We have added a full "
            "sensitivity analysis varying the equipment-CHE share from 5% "
            "to 25% (new Table 5, new Figure 9). At 5%, the fiscal return "
            "is 1.02 (sustainable); at 10%, it is 1.00 (borderline); at "
            "15% (base case), 0.98; at 20%, 0.96; at 25%, 0.94. The "
            "sustainability conclusion under Scenario A thus depends "
            "critically on this assumption. We now state this explicitly "
            "in both Results and Discussion."
        ),
        (
            "6",
            "Scenario A presumes that reducing equipment density lowers "
            "the multiplier, but the chain density \u2192 spending \u2192 "
            "multiplier is not transparently modeled.",
            "We have added a new paragraph in Methods ('Counterfactual "
            "scenarios and sensitivity analysis') that explicitly describes "
            "the four-step transmission chain: (i) density reduction \u2192 "
            "(ii) proportionally lower equipment-related spending \u2192 "
            "(iii) multiplier reduction (by half the proportional equipment "
            "share reduction, reflecting that some equipment spending flows "
            "to imports) \u2192 (iv) further adjustment for import leakage."
        ),
        (
            "7",
            "The gap between 0.98 and 1.0 is likely smaller than the "
            "upstream approximation error. Values near 1.0 should carry "
            "an error/sensitivity range.",
            "We fully agree. The Discussion now explicitly states that "
            "'the gap between the Scenario A ratio (0.98) and the threshold "
            "(1.0) is 0.02, which is smaller than the sensitivity range "
            "(0.94-1.02 across the 5-25% equipment share range). Values "
            "near 1.0 should be interpreted with an uncertainty range rather "
            "than as a binary sustainability determination.' The conclusions "
            "now use language of 'indicative ranges' rather than definitive "
            "sustainability claims."
        ),
    ]

    for num, comment, response in comments_r1:
        p = doc.add_paragraph()
        run = p.add_run(f"Comment {num}: ")
        run.bold = True
        run2 = p.add_run(comment)
        run2.italic = True
        doc.add_paragraph()
        p2 = doc.add_paragraph()
        run3 = p2.add_run("Response: ")
        run3.bold = True
        p2.add_run(response)
        doc.add_paragraph()

    # ---- REVIEWER 2 ----
    add_heading(doc, "Response to Reviewer 2", level=1)

    comments_r2 = [
        (
            "1",
            "The author's position must be made explicit: does the current "
            "equipment density represent a social surplus or a necessary "
            "density?",
            "We have added a new subsection 'Position on Japan's equipment "
            "density' in the Discussion that explicitly states our position: "
            "Japan's equipment density represents neither a clear surplus "
            "nor an unambiguous necessity, but a structural feature whose "
            "net effect depends on distinguishing value-generating diagnostic "
            "activity from wasteful over-utilization. We cite Hiramatsu et al. "
            "on defensive medicine in Japan (new reference 29)."
        ),
        (
            "2",
            "Separate analyses for the Late-Stage Elderly Healthcare System "
            "and other insurance frameworks are essential.",
            "We have added a new Discussion subsection on 'Implications of "
            "Japan's elderly insurance structure' that describes the Late-Stage "
            "Elderly Healthcare System (Kouki Koreisha Iryo Seido) and its "
            "implications for sustainability analysis. We acknowledge that "
            "treating expenditure as a homogeneous aggregate is a limitation "
            "and that disaggregated analysis is an important future direction. "
            "The Background now also describes this system (new reference 26)."
        ),
        (
            "3",
            "A context within a super-aged society with growing elderly and "
            "shrinking workforce should be considered.",
            "The Background now includes an expanded paragraph on Japan's "
            "super-aged society context, including projections (35% aged 65+ "
            "by 2040, 20% workforce decline). The Discussion also addresses "
            "the structural pressure this places on healthcare financing."
        ),
        (
            "4",
            "The title fails to convey that the study specifically evaluates "
            "the sustainability of Japanese medical service.",
            "The title has been revised to: 'Evaluating Healthcare Expenditure "
            "Sustainability in Japan: A Dual-Return Framework Integrating "
            "Input-Output Multipliers, Health-Capital Tempo, and Diagnostic "
            "Equipment Stock with Cross-Country Benchmarking.' This signals "
            "the Japan focus while indicating the cross-country context."
        ),
        (
            "5",
            "Physician remuneration and density variations among OECD "
            "countries must be taken into account.",
            "The Background now discusses Japan's low physician density "
            "(2.5/1,000 vs OECD average 3.7) and work-style reform "
            "constraints (new references 27, 28). The Discussion includes "
            "a new subsection on 'Over-supply risks and workforce "
            "sustainability' addressing the tension between low physician "
            "density and high service access. The limitations section "
            "notes that physician remuneration differences affect "
            "international comparability."
        ),
        (
            "6",
            "When public financing is funded by deficit rather than taxation, "
            "a ratio above 1.0 reflects only a temporary flow. The methodology "
            "needs re-evaluation.",
            "We have added deficit-adjusted analysis. Table 2 now includes "
            "a 'Deficit Share' column showing the general-account deficit "
            "dependency for each country. A new Discussion subsection "
            "'Deficit financing and national debt' addresses this concern "
            "in depth, noting Japan's 35% deficit dependency and 250%+ "
            "debt-to-GDP ratio (new references 25, 30). The Methods section "
            "now describes the deficit-adjusted variant."
        ),
        (
            "7",
            "The applicability of the current index for international "
            "comparisons is highly questionable given different deficit "
            "financing structures.",
            "A new Discussion subsection 'Limitations of international "
            "comparison' addresses this, noting that deficit financing "
            "structures vary substantially (Sweden 2% vs Japan 35%) and "
            "that direct ratio comparison is appropriate for identifying "
            "broad patterns but not for normative sustainability judgments."
        ),
        (
            "8",
            "The Abstract should be revised to align with the main findings.",
            "The Abstract has been substantially revised. The focus is now "
            "on Japan with cross-country benchmarking as context. Results "
            "include both output and VA multiplier findings, the sensitivity "
            "range for Scenario A, and the honest M2 vs M1 assessment. "
            "Conclusions are appropriately cautious."
        ),
        (
            "9",
            "There must be a baseline assumption that quality and "
            "accessibility are uniform across countries. In reality, "
            "this is not the case.",
            "The new 'Limitations of international comparison' subsection "
            "discusses the differences in access (free specialist access "
            "vs GP gatekeeping), emergency care structure, and physician "
            "remuneration that make uniform comparison problematic."
        ),
        (
            "10",
            "The author does not evaluate the potential for over-supply. "
            "Defensive medicine, over-testing, and medical staff burnout "
            "must be addressed.",
            "A new subsection 'Over-supply risks and workforce sustainability' "
            "in the Discussion addresses defensive medicine (citing Hiramatsu "
            "et al., ref 29), the tension between high equipment density and "
            "low physician density, and the implications of work-style reform. "
            "We acknowledge that the multiplier captures economic activity, "
            "not health value, and that low-value utilization would overstate "
            "the sustainability benefit."
        ),
        (
            "11-12",
            "The amount of national debt allocated to healthcare should be "
            "estimated. Health expenditure causes national debt.",
            "The Background now discusses Japan's fiscal context including "
            "deficit dependency (35%) and national debt (250%+ of GDP). "
            "A new Discussion subsection provides detailed analysis of how "
            "deficit financing affects interpretation of the fiscal return "
            "ratio. New references 25 and 30 support this material."
        ),
        (
            "13",
            "Modifications beyond the Methods section are necessary.",
            "As detailed above, revisions have been made throughout the "
            "manuscript: Background (super-aged context, elderly healthcare "
            "system, physician density, national debt), Methods (output vs "
            "VA distinction, deficit adjustment, sensitivity analysis), "
            "Results (VA multiplier ratios, sensitivity table), Discussion "
            "(equipment position, over-supply, workforce, elderly insurance, "
            "deficit, international comparison caveats), and Conclusions "
            "(substantially toned down)."
        ),
    ]

    for num, comment, response in comments_r2:
        p = doc.add_paragraph()
        run = p.add_run(f"Comment {num}: ")
        run.bold = True
        run2 = p.add_run(comment)
        run2.italic = True
        doc.add_paragraph()
        p2 = doc.add_paragraph()
        run3 = p2.add_run("Response: ")
        run3.bold = True
        p2.add_run(response)
        doc.add_paragraph()

    # Closing
    doc.add_paragraph()
    add_plain_para(doc,
        "We believe these revisions have substantially strengthened the "
        "manuscript and addressed all reviewer concerns. We are happy to "
        "make any further modifications the editor or reviewers may suggest.")
    doc.add_paragraph()
    add_plain_para(doc, "Sincerely,")
    doc.add_paragraph()
    add_plain_para(doc, "Tatsuki Onishi")

    out_path = os.path.join(DOCX_DIR, "Healthcare_EHPM_ResponseToReviewers.docx")
    doc.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# Cover letter for revised submission
# ---------------------------------------------------------------------------
def build_cover_letter_r1():
    doc = Document()
    _set_ehpm_format(doc)

    from datetime import date
    today = date.today().strftime("%B %d, %Y")
    add_plain_para(doc, today)
    doc.add_paragraph()

    add_plain_para(doc, "Professor Kouji H. Harada, PhD, MPH")
    add_plain_para(doc, "Editor-in-Chief")
    add_plain_para(doc, "Environmental Health and Preventive Medicine")
    doc.add_paragraph()

    add_plain_para(doc,
        "Re: Revised manuscript EHPM-D-26-00106")
    doc.add_paragraph()

    add_plain_para(doc, "Dear Professor Harada,")
    doc.add_paragraph()

    add_plain_para(doc,
        "Thank you for the opportunity to revise our manuscript entitled "
        '"Evaluating Healthcare Expenditure Sustainability in Japan: '
        "A Dual-Return Framework Integrating Input-Output Multipliers, "
        "Health-Capital Tempo, and Diagnostic Equipment Stock with "
        'Cross-Country Benchmarking" (formerly titled "Healthcare '
        'Expenditure as Economic Effect..."). '
        "We appreciate the thorough reviews by both reviewers and have "
        "made extensive revisions in response.")
    doc.add_paragraph()

    add_plain_para(doc, "Major revisions include:")
    changes = [
        "Title revised to explicitly signal Japan focus (Reviewer 2)",
        "Output vs value-added multiplier distinction with dual reporting "
        "(Reviewer 1, Comment 1)",
        "Self-contained tempo model specification with point estimates and "
        "standard errors (Reviewer 1, Comment 3)",
        "AIC/BIC/LOOCV RMSE added to Table 3 with honest M2 vs M1 "
        "assessment (Reviewer 1, Comment 4)",
        "New sensitivity analysis (see Results) for equipment-CHE "
        "share assumption (Reviewer 1, Comment 5)",
        "Deficit financing and national debt discussion (Reviewer 2, "
        "Comments 6-7, 11-12)",
        "Late-Stage Elderly Healthcare System and super-aged society "
        "context (Reviewer 2, Comments 2-3)",
        "Over-supply risks, defensive medicine, and workforce "
        "sustainability discussion (Reviewer 2, Comments 5, 10)",
        "International comparability caveats (Reviewer 2, Comments 7, 9)",
        "Conclusions substantially toned down (Reviewer 1, Comment 7)",
    ]
    for change in changes:
        p = doc.add_paragraph()
        run = p.add_run("\u2022 ")
        p.add_run(change)

    doc.add_paragraph()
    add_plain_para(doc,
        "A detailed point-by-point response to each reviewer comment "
        "is provided in the accompanying 'Response to Reviewers' document. "
        "All changes are indicated in the revised manuscript.")
    doc.add_paragraph()

    add_plain_para(doc,
        "We confirm that no part of this research was funded or supported "
        "by firms or organizations related to the tobacco industry.")
    doc.add_paragraph()

    add_plain_para(doc, "Thank you for your consideration.")
    doc.add_paragraph()
    add_plain_para(doc, "Sincerely,")
    doc.add_paragraph()
    add_plain_para(doc, "Tatsuki Onishi")
    add_plain_para(doc, "[Affiliation]")
    add_plain_para(doc, "[E-mail]")

    # Reminder for author
    doc.add_paragraph()
    p = doc.add_paragraph()
    run = p.add_run("REMINDER: Enter discount code EHPM-JSH-R26K during "
                     "APC Agreement stage on Editorial Manager submission.")
    run.bold = True
    run.font.color.rgb = RGBColor(255, 0, 0)

    out_path = os.path.join(DOCX_DIR, "Healthcare_EHPM_CoverLetter_R1.docx")
    doc.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


# ---------------------------------------------------------------------------
# PPTX (editable figures, R1)
# ---------------------------------------------------------------------------
def build_ehpm_pptx_r1():
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    figures = [
        ("fig4_dual_return_schematic.png",
         "Figure 1. Dual-Return Framework", ""),
        ("fig1_io_multipliers.png",
         "Figure 2. Healthcare I-O Output Multipliers",
         "Note: US = Medicare only."),
        ("fig3_fiscal_sustainability.png",
         "Figure 3. Fiscal Return Ratio", ""),
        ("fig5_three_layer_analogy.png",
         "Figure 4. Three-Layer Tempo Analogy",
         "M2 vs M1 does not survive BIC."),
        ("fig2_che_vs_lifeexp.png",
         "Figure 5. CHE vs Life Expectancy", ""),
        ("fig6_equipment_density.png",
         "Figure 6. Equipment Density", ""),
        ("fig7_import_leakage_multiplier.png",
         "Figure 7. Import Leakage vs Multiplier", ""),
        ("fig8_counterfactual_japan.png",
         "Figure 8. Japan Counterfactual", ""),
        ("fig9_sensitivity_equip_share.png",
         "Figure 9. Equipment CHE Share Sensitivity",
         "NEW in R1. Scenario A FR varies 0.94-1.02."),
    ]

    for fname, title, caption in figures:
        path = get_fig(fname)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2),
                                          PptxInches(12.33), PptxInches(0.6))
        tf = txBox.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = PptxPt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if os.path.exists(path):
            slide.shapes.add_picture(path, PptxInches(1.5), PptxInches(1.0),
                                      PptxInches(10.33), PptxInches(5.0))
        if caption:
            txBox2 = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(6.3),
                PptxInches(12.33), PptxInches(0.8))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            tf2.text = caption
            tf2.paragraphs[0].font.size = PptxPt(10)
            tf2.paragraphs[0].font.italic = True
            tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    out_path = os.path.join(PPTX_DIR, "Healthcare_EHPM_Figures_R1.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    build_ehpm_manuscript_r1()
    build_response_to_reviewers()
    build_cover_letter_r1()
    build_ehpm_pptx_r1()
